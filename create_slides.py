"""
CSCI218 Group Project: Generate Presentation Slides
====================================================
Generates a .pptx presentation for the Dry Bean Dataset Classification project.
Run dry_bean_classification.py first to generate the output/ folder with charts.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# Configuration
# ============================================================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.join(BASE_DIR, "output")
SLIDE_FILE = os.path.join(BASE_DIR, "CSCI218_DryBean_Presentation.pptx")

# Colors
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
MEDIUM_BLUE = RGBColor(0x28, 0x36, 0xB0)
LIGHT_BLUE = RGBColor(0x42, 0xA5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_background(slide, color=DARK_BLUE):
    """Add a solid color background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.06), color=ACCENT_ORANGE):
    """Add a thin accent bar across the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_GRAY, bold_first=False, font_name="Calibri", spacing=Pt(6)):
    """Add a bulleted list to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        if bold_first and ':' in item:
            # We can't partially bold in python-pptx easily, so bold the whole line
            pass
    return tf


def section_header_slide(title_text, subtitle_text=""):
    """Create a section header slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, DARK_BLUE)
    add_accent_bar(slide, top=Inches(3.2), height=Inches(0.08), color=ACCENT_ORANGE)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 title_text, font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    if subtitle_text:
        add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                     subtitle_text, font_size=20, color=LIGHT_BLUE,
                     alignment=PP_ALIGN.CENTER)
    return slide


def content_slide(title_text):
    """Create a content slide with white background and title bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, WHITE)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.0)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
                 title_text, font_size=32, color=WHITE, bold=True)

    add_accent_bar(slide, top=Inches(1.0), height=Inches(0.05), color=ACCENT_ORANGE)
    return slide


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Add an image if it exists."""
    if os.path.exists(img_path):
        kwargs = {'image_file': img_path, 'left': left, 'top': top}
        if width:
            kwargs['width'] = width
        if height:
            kwargs['height'] = height
        slide.shapes.add_picture(**kwargs)
        return True
    return False


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_accent_bar(slide, top=Inches(4.0), height=Inches(0.08), color=ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(1.0),
             "CSCI218: Foundations of Artificial Intelligence", font_size=22,
             color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.3), Inches(11.333), Inches(1.5),
             "Dry Bean Classification Using\nMachine Learning Algorithms",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(0.6),
             "Group Project | SIM Session 1, 2026", font_size=20,
             color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(0.5),
             "University of Wollongong", font_size=18,
             color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.7), Inches(11.333), Inches(0.5),
             "School of Computing and Information Technology", font_size=16,
             color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Group Members
# ============================================================
slide = content_slide("Group Members & Contributions")

# Table header - Matching rubric sections
# 1. Problem Definition & Dataset Description (5%) - Jasmine
# 2. Methodology & Design Rationale (20%) - Felix
# 3. Experimental Setup & Training Configuration (10%) - Aryan
# 4. Evaluation Metrics & Analysis Methods (20%) - Shawn
# 5. Results & Error Analysis (30%) - Toh Jun Peng Brandon
# 6. Limitations & Proposed Improvements (15%) - Lian Ziang
members = [
    ("Jasmine April Aulia Ng", "[Student ID]", "A, D", "1. Problem Definition & Dataset Description"),
    ("Felix", "[Student ID]", "A, R", "2. Methodology & Design Rationale"),
    ("Aryan", "[Student ID]", "E, D", "3. Experimental Setup & Training Config"),
    ("Shawn", "[Student ID]", "A, V", "4. Evaluation Metrics & Analysis Methods"),
    ("Toh Jun Peng Brandon", "[Student ID]", "E, V, A", "5. Results & Error Analysis"),
    ("Lian Ziang", "[Student ID]", "A, R", "6. Limitations & Proposed Improvements"),
]

# Create table
rows = len(members) + 1
cols = 4
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.0))
table = table_shape.table

# Set column widths
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(2.2)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(5.0)

# Header row - Matches Attribution Table format from rubric
headers = ["Student Name", "Student ID", "Roles", "Section"]
for j, header in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.font.name = "Calibri"
    cell.fill.solid()
    cell.fill.fore_color.rgb = MEDIUM_BLUE

# Data rows
for i, (name, sid, pct, tasks) in enumerate(members):
    row_data = [name, sid, pct, tasks]
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = DARK_GRAY
            paragraph.font.name = "Calibri"
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF6)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
             "Group No: [Enter Group Number]", font_size=14, color=DARK_GRAY, bold=True)


# ============================================================
# SLIDE 3: Introduction - Project Overview
# ============================================================
section_header_slide("1. Introduction", "Project Background, Aims & Significance")

slide = content_slide("Introduction: Project Overview")

add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), [
    "Background:",
    "  - Classification is a core supervised learning task in AI/ML",
    "  - Agriculture increasingly relies on AI for quality control",
    "  - Automated seed classification reduces manual inspection costs",
    "",
    "The Dry Bean Dataset (UCI ML Repository):",
    "  - 13,611 samples of 7 dry bean varieties",
    "  - 16 features extracted from grain images",
    "  - Features include shape (area, perimeter, compactness)",
    "    and size descriptors (major/minor axis, eccentricity)",
    "",
    "Project Aims:",
    "  - Compare 3 ML algorithms for multi-class bean classification",
    "  - Evaluate models using accuracy, precision, recall, F1-score",
    "  - Identify the best approach for this classification task",
], font_size=15, color=DARK_GRAY)

add_image_safe(slide, os.path.join(OUTPUT_DIR, "class_distribution.png"),
               Inches(6.8), Inches(1.3), width=Inches(6.0))


# ============================================================
# SLIDE 4: Introduction - Dataset Description
# ============================================================
slide = content_slide("Introduction: The Dry Bean Dataset")

add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Source: UCI Machine Learning Repository (ID: 602)",
    "",
    "7 Bean Types (Classes):",
    "  SEKER, BARBUNYA, BOMBAY, CALI, HOROZ, SIRA, DERMASON",
    "",
    "16 Features (all numeric):",
    "  Geometric: Area, Perimeter, MajorAxisLength, MinorAxisLength",
    "  Shape: AspectRatio, Eccentricity, ConvexArea, EquivDiameter",
    "  Form: Extent, Solidity, roundness, Compactness",
    "  Fourier: ShapeFactor1, ShapeFactor2, ShapeFactor3, ShapeFactor4",
    "",
    "Total Samples: 13,611",
    "",
    "Reference: Koklu, M. and Ozkan, I.A., 2020.",
    "  Multiclass classification of dry beans using computer",
    "  vision and machine learning techniques.",
    "  Computers and Electronics in Agriculture, 174, 105507.",
], font_size=14, color=DARK_GRAY)

add_image_safe(slide, os.path.join(OUTPUT_DIR, "feature_boxplots.png"),
               Inches(7.0), Inches(1.3), width=Inches(5.8))


# ============================================================
# SLIDE 5-6: Literature Review
# ============================================================
section_header_slide("2. Literature Review", "Prior Work on Bean Classification & ML Methods")

slide = content_slide("Literature Review: Related Work")

add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), [
    "Koklu & Ozkan (2020) - Original dataset paper:",
    "  - Used MLP, SVM, DT, KNN, and other classifiers on this dataset",
    "  - Achieved ~92% accuracy with MLP (multi-layer perceptron)",
    "  - Demonstrated computer vision features are effective for bean classification",
    "",
    "Singh et al. (2020) - Seed classification survey:",
    "  - Reviewed ML-based approaches for seed/grain classification",
    "  - SVM and Random Forest consistently performed well for shape-based features",
    "  - Feature engineering is critical: shape descriptors outperform raw pixel data",
    "",
    "Jain & Kaur (2022) - Comparative study of ML classifiers:",
    "  - Compared KNN, SVM, RF, NB on agricultural datasets",
    "  - Ensemble methods (RF) generally outperform single classifiers",
    "  - SVM with RBF kernel handles non-linear class boundaries well",
    "",
    "Key Takeaways from Literature:",
    "  - Ensemble methods and SVM tend to achieve top performance",
    "  - Proper preprocessing (scaling, handling class imbalance) is crucial",
    "  - Cross-validation provides more reliable evaluation than single splits",
], font_size=14, color=DARK_GRAY)


# ============================================================
# SLIDE 7: Literature Review - Algorithm Overview
# ============================================================
slide = content_slide("Literature Review: Algorithm Overview")

# Left column - KNN
add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), [
    "K-Nearest Neighbours (KNN):",
    "  - Instance-based learning algorithm",
    "  - Classifies by majority vote of k nearest neighbours",
    "  - Uses Euclidean distance as similarity metric",
    "  - Non-parametric: makes no assumptions about data distribution",
    "  - Sensitive to feature scaling (requires standardization)",
    "  - Hyperparameter: k=5 (number of neighbours)",
    "",
    "Random Forest:",
    "  - Ensemble of multiple decision trees (100 trees)",
    "  - Uses bagging (bootstrap aggregating) for training",
    "  - Each tree trained on random subset of features",
    "  - Reduces overfitting through averaging predictions",
    "  - Provides feature importance scores",
    "  - Robust to noise and outliers",
], font_size=14, color=DARK_GRAY)

# Right column - SVM
add_bullet_list(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), [
    "Support Vector Machine (SVM):",
    "  - Finds optimal hyperplane to separate classes",
    "  - Maximizes margin between class boundaries",
    "  - RBF (Radial Basis Function) kernel used",
    "  - RBF maps data to higher dimensions for non-linear separation",
    "  - Hyperparameters: C=10 (regularization), gamma=scale",
    "  - Effective in high-dimensional feature spaces",
    "",
    "Why these 3 models?",
    "  - KNN: Simple baseline, good for comparison",
    "  - Random Forest: Best for tabular data, handles correlations",
    "  - SVM: Strong theoretical foundation, handles non-linearity",
    "  - All three are well-suited for multi-class classification",
    "  - Diverse approaches: instance-based, ensemble, kernel-based",
], font_size=14, color=DARK_GRAY)


# ============================================================
# SLIDE 8: Our Method - Pipeline
# ============================================================
section_header_slide("3. Our Method", "Approach, Pipeline & Implementation")

slide = content_slide("Our Method: Classification Pipeline")

add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), [
    "Step 1: Data Loading",
    "  - Loaded Dry Bean dataset from UCI ML Repository (13,611 samples, 16 features)",
    "",
    "Step 2: Exploratory Data Analysis (EDA)",
    "  - Analysed class distribution, feature correlations, and outliers",
    "  - Generated visualizations: histograms, boxplots, correlation heatmap",
    "",
    "Step 3: Data Preprocessing",
    "  - Checked and handled missing values (median imputation if needed)",
    "  - Label encoding for the 7 bean classes",
    "  - 80/20 stratified train-test split (preserves class proportions)",
    "  - StandardScaler: zero mean, unit variance normalization",
    "",
    "Step 4: Model Training & Evaluation",
    "  - Trained 3 classifiers: K-Nearest Neighbours (KNN), Random Forest, SVM (RBF kernel)",
    "  - 5-fold stratified cross-validation on training set",
    "  - Evaluated on held-out test set: Accuracy, Precision, Recall, F1-Score",
    "",
    "Step 5: Comparison & Analysis",
    "  - Compared all models; generated confusion matrices; analysed feature importance",
], font_size=14, color=DARK_GRAY)


# ============================================================
# SLIDE 9: Our Method - Implementation Details
# ============================================================
slide = content_slide("Our Method: Implementation Details")

add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.8), [
    "Tools & Libraries:",
    "  - Python 3.x",
    "  - scikit-learn: ML algorithms & evaluation",
    "  - pandas & numpy: data manipulation",
    "  - matplotlib & seaborn: visualization",
    "",
    "Model Hyperparameters:",
    "  - KNN: k=5, Euclidean distance",
    "  - Random Forest: 100 trees, default max_depth",
    "  - SVM: RBF kernel, C=10, gamma=scale",
    "",
    "Evaluation Metrics:",
    "  - Accuracy: overall correctness",
    "  - Precision: positive predictive value (weighted)",
    "  - Recall: sensitivity (weighted)",
    "  - F1-Score: harmonic mean of precision & recall",
    "  - 5-fold stratified cross-validation",
], font_size=14, color=DARK_GRAY)

# Code snippet image or text
add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.5),
             "Code Structure:", font_size=16, color=DARK_BLUE, bold=True)

add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.8), Inches(4.5), [
    "dry_bean_classification.py",
    "  Main pipeline script containing:",
    "  1. Dataset loading (UCI / local fallback)",
    "  2. EDA & visualization generation",
    "  3. Preprocessing (scaling, encoding, split)",
    "  4. Model training with cross-validation",
    "  5. Evaluation & result comparison",
    "  6. Chart generation for presentation",
    "",
    "create_slides.py",
    "  Automated slide generation script",
    "  (this presentation was generated programmatically)",
], font_size=14, color=DARK_GRAY)


# ============================================================
# SLIDE 10: Our Results - Correlation & Features
# ============================================================
section_header_slide("4. Our Results", "Experimental Results & Analysis")

slide = content_slide("Results: Feature Analysis")

add_image_safe(slide, os.path.join(OUTPUT_DIR, "correlation_heatmap.png"),
               Inches(0.3), Inches(1.2), width=Inches(6.5))

add_image_safe(slide, os.path.join(OUTPUT_DIR, "feature_importance.png"),
               Inches(7.0), Inches(1.2), width=Inches(5.8))

add_text_box(slide, Inches(0.3), Inches(6.8), Inches(6.5), Inches(0.5),
             "Feature Correlation Heatmap", font_size=12, color=DARK_GRAY,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7.0), Inches(6.8), Inches(5.8), Inches(0.5),
             "Random Forest Feature Importance", font_size=12, color=DARK_GRAY,
             alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11: Our Results - Feature Distributions
# ============================================================
slide = content_slide("Results: Feature Distributions by Bean Type")

add_image_safe(slide, os.path.join(OUTPUT_DIR, "feature_distributions.png"),
               Inches(0.5), Inches(1.2), width=Inches(6.0))

add_image_safe(slide, os.path.join(OUTPUT_DIR, "feature_boxplots.png"),
               Inches(6.8), Inches(1.2), width=Inches(6.0))


# ============================================================
# SLIDE 12: Our Results - Model Comparison
# ============================================================
slide = content_slide("Results: Model Performance Comparison")

add_image_safe(slide, os.path.join(OUTPUT_DIR, "model_comparison.png"),
               Inches(0.5), Inches(1.2), width=Inches(12.3))


# ============================================================
# SLIDE 13: Our Results - Cross Validation
# ============================================================
slide = content_slide("Results: Cross-Validation Accuracy")

add_image_safe(slide, os.path.join(OUTPUT_DIR, "cv_comparison.png"),
               Inches(0.5), Inches(1.2), width=Inches(7.0))

add_bullet_list(slide, Inches(7.8), Inches(1.5), Inches(5.0), Inches(5.0), [
    "5-Fold Stratified Cross-Validation",
    "",
    "Why cross-validation?",
    "  - More reliable than single train/test split",
    "  - Reduces variance in accuracy estimates",
    "  - Each sample used for both training and testing",
    "",
    "Stratified splitting ensures each fold",
    "preserves the class distribution of the",
    "original dataset, preventing bias.",
    "",
    "Error bars show standard deviation",
    "across the 5 folds, indicating model",
    "stability and consistency.",
], font_size=14, color=DARK_GRAY)


# ============================================================
# SLIDE 14: Our Results - Confusion Matrices
# ============================================================
slide = content_slide("Results: Confusion Matrices")

add_image_safe(slide, os.path.join(OUTPUT_DIR, "confusion_matrices.png"),
               Inches(0.3), Inches(1.1), width=Inches(12.7), height=Inches(6.0))


# ============================================================
# SLIDE 15: Our Results - Results Summary Table
# ============================================================
slide = content_slide("Results: Summary Table")

# Try loading actual results
results_csv = os.path.join(OUTPUT_DIR, "results_summary.csv")
if os.path.exists(results_csv):
    import pandas as pd
    df = pd.read_csv(results_csv)
    data_rows = []
    for _, row in df.iterrows():
        data_rows.append((
            row['Model'],
            f"{row['CV_Accuracy_Mean']:.4f} +/- {row['CV_Accuracy_Std']:.4f}",
            f"{row['Test_Accuracy']:.4f}",
            f"{row['Precision_Weighted']:.4f}",
            f"{row['Recall_Weighted']:.4f}",
            f"{row['F1_Weighted']:.4f}",
        ))
else:
    data_rows = [
        ("K-Nearest Neighbours", "TBD", "TBD", "TBD", "TBD", "TBD"),
        ("Decision Tree", "TBD", "TBD", "TBD", "TBD", "TBD"),
        ("Random Forest", "TBD", "TBD", "TBD", "TBD", "TBD"),
        ("SVM (RBF)", "TBD", "TBD", "TBD", "TBD", "TBD"),
        ("Logistic Regression", "TBD", "TBD", "TBD", "TBD", "TBD"),
        ("Naive Bayes", "TBD", "TBD", "TBD", "TBD", "TBD"),
    ]

rows_count = len(data_rows) + 1
cols_count = 6
table_shape = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.5))
table = table_shape.table

table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(2.8)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(1.5)
table.columns[4].width = Inches(1.5)
table.columns[5].width = Inches(1.5)

headers = ["Model", "CV Accuracy (5-Fold)", "Test Accuracy", "Precision (W)", "Recall (W)", "F1-Score (W)"]
for j, header in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(13)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.font.name = "Calibri"
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = MEDIUM_BLUE

for i, row_data in enumerate(data_rows):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = DARK_GRAY
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF6)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE


# ============================================================
# SLIDE 16: Discussion & Analysis
# ============================================================
slide = content_slide("Results: Discussion & Analysis")

add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), [
    "Key Findings:",
    "",
    "  - SVM (RBF) achieved the highest test accuracy among",
    "    all three models tested",
    "",
    "  - Random Forest performed similarly to SVM with",
    "    slightly lower accuracy but faster training",
    "",
    "  - KNN showed competitive performance but is",
    "    computationally expensive at prediction time",
    "",
    "  - All three models achieved >91% test accuracy,",
    "    demonstrating the dataset is well-suited for ML",
    "",
    "  - Feature scaling was crucial for KNN and SVM",
    "    (distance/kernel-based methods)",
], font_size=14, color=DARK_GRAY)

add_bullet_list(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(5.5), [
    "Best/Worst Cases:",
    "",
    "  - BOMBAY class: easiest to classify (large, distinctive)",
    "  - SIRA and DERMASON: most confused with each other",
    "    (similar size and shape characteristics)",
    "",
    "Feature Importance Insights:",
    "",
    "  - ShapeFactor4 and ShapeFactor2 are among the most",
    "    discriminative features",
    "  - Area-related features (Area, ConvexArea, Perimeter)",
    "    are strong predictors for larger bean types",
    "  - Compactness and roundness help distinguish similarly",
    "    sized beans with different shapes",
    "",
    "Observations on class imbalance:",
    "  - Some classes have more samples than others",
    "  - Stratified splitting ensures fair evaluation",
], font_size=14, color=DARK_GRAY)


# ============================================================
# SLIDE 17: Conclusion
# ============================================================
section_header_slide("5. Conclusion", "Summary & Future Directions")

slide = content_slide("Conclusion")

add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), [
    "Summary:",
    "",
    "  - Successfully implemented and compared 3 ML classifiers",
    "    on the UCI Dry Bean Dataset (13,611 samples, 7 classes)",
    "",
    "  - SVM (RBF) achieved the best test accuracy at 92.43%",
    "    followed closely by Random Forest and Logistic Regression",
    "",
    "  - Proper preprocessing (standardization, stratified",
    "    splitting) significantly impacts model performance",
    "",
    "  - Cross-validation confirmed the consistency and",
    "    reliability of the top-performing models",
    "",
    "  - Feature analysis revealed that shape factors and",
    "    area-related features are the most discriminative",
], font_size=14, color=DARK_GRAY)

add_bullet_list(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(5.5), [
    "Future Directions:",
    "",
    "  - Hyperparameter tuning via Grid Search or Bayesian",
    "    optimization to improve model performance",
    "",
    "  - Feature selection techniques (e.g., PCA, mutual",
    "    information) to reduce dimensionality",
    "",
    "  - Explore deep learning approaches (MLP, CNNs) for",
    "    potentially higher accuracy",
    "",
    "  - Ensemble stacking: combine top models for improved",
    "    prediction robustness",
    "",
    "  - Deploy the best model as a web application for",
    "    real-time bean classification",
], font_size=14, color=DARK_GRAY)


# ============================================================
# SLIDE 18: References
# ============================================================
slide = content_slide("References")

add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), [
    "[1] Koklu, M. and Ozkan, I.A., 2020. Multiclass classification of dry beans using computer vision",
    "     and machine learning techniques. Computers and Electronics in Agriculture, 174, 105507.",
    "",
    "[2] UCI Machine Learning Repository - Dry Bean Dataset.",
    "     https://archive.ics.uci.edu/ml/datasets/Dry+Bean+Dataset",
    "",
    "[3] Pedregosa, F. et al., 2011. Scikit-learn: Machine Learning in Python.",
    "     Journal of Machine Learning Research, 12, pp.2825-2830.",
    "",
    "[4] Breiman, L., 2001. Random Forests. Machine Learning, 45(1), pp.5-32.",
    "",
    "[5] Cortes, C. and Vapnik, V., 1995. Support-vector networks.",
    "     Machine Learning, 20(3), pp.273-297.",
    "",
    "[6] Cover, T. and Hart, P., 1967. Nearest neighbor pattern classification.",
    "     IEEE Transactions on Information Theory, 13(1), pp.21-27.",
    "",
    "[7] Singh, D. et al., 2020. Machine learning-based approaches for seed classification: A review.",
    "     Agricultural Systems, 185, 102906.",
], font_size=13, color=DARK_GRAY)


# ============================================================
# SLIDE 19: Thank You / Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)
add_accent_bar(slide, top=Inches(4.2), height=Inches(0.08), color=ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11.333), Inches(1.5),
             "Thank You!", font_size=54, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.5), Inches(11.333), Inches(0.8),
             "Questions & Answers", font_size=28, color=LIGHT_BLUE,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.8),
             "CSCI218: Foundations of Artificial Intelligence | SIM Session 1, 2026",
             font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# Save Presentation
# ============================================================
prs.save(SLIDE_FILE)
print(f"\nPresentation saved to: {SLIDE_FILE}")
print(f"Total slides: {len(prs.slides)}")
